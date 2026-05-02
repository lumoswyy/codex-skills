"""
图文页 - 终版：溢出保护 + 填满 + 左右/上下自动布局
"""
from .base import BaseTemplate, hex_to_rgb
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree
import re, os


def _set_lnsp(tf, pt_val):
    for para in tf.paragraphs:
        pPr = para._p.get_or_add_pPr()
        el = pPr.find(qn("a:lnSpc"))
        if el is None: el = etree.SubElement(pPr, qn("a:lnSpc"))
        sp = el.find(qn("a:spcPts"))
        if sp is None: sp = etree.SubElement(el, qn("a:spcPts"))
        sp.set("val", str(int(pt_val * 100)))


def _aspect(path):
    try:
        from PIL import Image as PILImage
        with PILImage.open(path) as im:
            w, h = im.size; return w / h if h else 1.5
    except Exception: return 1.5


def _calc(texts, avail_h, cw):
    """字号从大到小，保证 sum(heights) <= avail_h"""
    for fs in range(20, 9, -1):
        cpl = max(1, int((cw - 0.15) * 72 / (fs * 0.58)))
        heights = [fs / 72 * 1.28 * max(1, -(-len(t) // cpl)) + 0.05 for t in texts]
        if sum(heights) <= avail_h:
            return fs, heights
    fs = 9
    cpl = max(1, int((cw - 0.15) * 72 / (fs * 0.58)))
    heights = [fs / 72 * 1.28 * max(1, -(-len(t) // cpl)) + 0.05 for t in texts]
    return fs, heights


class ContentFigureTemplate(BaseTemplate):
    HEADER_Y = 0.68
    FOOT_Y   = 7.28

    def render(self, data):
        self.draw_nav_sidebar()
        cx = self.CONTENT_X
        cw = self.W - cx - 0.22

        title    = getattr(data, "title",    "")
        points   = getattr(data, "points",   []) or []
        figure   = getattr(data, "figure",   None)
        part_num = getattr(data, "part_num", "")

        # 多图列表：优先用 figures 字段，兜底用 figure
        figures  = getattr(data, "figures",  []) or []
        if not figures and figure:
            figures = [figure]

        self.draw_header_bar(part_num=part_num, title=title)

        top   = self.HEADER_Y + 0.12
        avail = self.FOOT_Y - top

        n_figs = len(figures)

        if n_figs == 0:
            self._text_full(points, cx, cw, top, avail)
        elif n_figs == 1:
            # 原有单图逻辑
            fig  = figures[0]
            path = getattr(fig, "path", "") or ""
            cap  = getattr(fig, "caption", getattr(fig, "label", "图"))
            asp  = _aspect(path) if (path and os.path.isfile(path)) else 1.6
            if not points:
                self._img_centered(path, cap, cx, top, cw, avail, asp)
            elif asp > 1.75:
                self._top_bottom(points, path, cap, cx, top, cw, avail, asp)
            else:
                self._left_right(points, path, cap, cx, top, cw, avail, asp)
        elif n_figs == 2:
            self._two_figs(points, figures, cx, top, cw, avail)
        elif n_figs == 3:
            self._three_figs(points, figures, cx, top, cw, avail)
        else:
            self._four_figs(points, figures, cx, top, cw, avail)

    # ── 上文下图 ─────────────────────────────────────────────

    def _top_bottom(self, points, path, caption, cx, top, cw, avail, asp):
        CAP_H  = 0.24
        texts  = [str(p).strip() for p in points]
        n      = len(texts)

        # 先给图片分配 55% 高度，文字用剩余
        img_avail = avail * 0.55 - CAP_H
        txt_avail = avail * 0.45 - 0.12

        fs, heights = _calc(texts, txt_avail, cw)
        total  = sum(heights)
        gap    = max(0, txt_avail - total) / max(n, 1)

        y = top
        y = self._draw_texts(texts, cx, cw, y, fs, heights, gap)

        sep_y = y + 0.05
        self.add_line(cx, sep_y, cw, color="CCCCCC", width_pt=0.5)

        img_y = sep_y + 0.07
        remaining_h = self.FOOT_Y - img_y - CAP_H
        img_w, img_h = self._fit(cw, max(remaining_h, 0.5), asp)
        ox = cx + (cw - img_w) / 2
        self._draw_img(path, ox, img_y, img_w, img_h)
        self._caption(caption, cx, img_y + img_h + 0.03, cw, CAP_H)

    # ── 左文右图 ─────────────────────────────────────────────

    def _left_right(self, points, path, caption, cx, top, cw, avail, asp):
        CAP_H  = 0.22
        texts  = [str(p).strip() for p in points]
        n      = len(texts)

        # 图片高 = avail - 图注，宽按比例，最多占内容区 50%
        img_h = avail - CAP_H
        img_w = img_h * asp
        if img_w > cw * 0.50:
            img_w = cw * 0.50
            img_h = img_w / asp

        txt_w = cw - img_w - 0.18
        img_x = cx + txt_w + 0.18
        img_y = top + max(0, (avail - img_h - CAP_H) / 2)

        fs, heights = _calc(texts, avail, txt_w)
        total  = sum(heights)
        gap    = max(0, avail - total) / max(n, 1)
        self._draw_texts(texts, cx, txt_w, top, fs, heights, gap)

        # 竖分隔线
        self.slide.shapes.add_connector(
            1,
            Inches(img_x - 0.10), Inches(top + 0.1),
            Inches(img_x - 0.10), Inches(top + avail - 0.1)
        ).line.color.rgb = hex_to_rgb("CCCCCC")

        self._draw_img(path, img_x, img_y, img_w, img_h)
        self._caption(caption, img_x, img_y + img_h + 0.03, img_w, CAP_H)

    # ── 辅助 ─────────────────────────────────────────────────

    def _draw_texts(self, texts, cx, cw, y_start, fs, heights, gap):
        y      = y_start
        lnsp   = fs * 1.28
        for i, text in enumerate(texts):
            box_h = heights[i] + gap
            tb = self.slide.shapes.add_textbox(
                Inches(cx + 0.05), Inches(y),
                Inches(cw - 0.05), Inches(box_h))
            tf = tb.text_frame; tf.word_wrap = True
            _set_lnsp(tf, lnsp)
            p = tf.paragraphs[0]
            r0 = p.add_run(); r0.text = "➢  "
            r0.font.name = self.FONT_EN; r0.font.size = Pt(fs)
            r0.font.color.rgb = hex_to_rgb(self.COLOR_PRIMARY)
            self._rich(p, text, fs)
            y += box_h
        return y

    def _fit(self, max_w, max_h, asp):
        w = max_w; h = w / asp
        if h > max_h: h = max_h; w = h * asp
        if w > max_w: w = max_w; h = w / asp
        return w, h

    def _img_centered(self, path, caption, cx, top, cw, avail, asp):
        CAP_H = 0.24
        w, h  = self._fit(cw, avail - CAP_H, asp)
        ox    = cx + (cw - w) / 2
        oy    = top + max(0, (avail - h - CAP_H) / 2)
        self._draw_img(path, ox, oy, w, h)
        self._caption(caption, cx, oy + h + 0.03, cw, CAP_H)

    def _draw_img(self, path, x, y, w, h):
        if path and os.path.isfile(path):
            try:
                self.slide.shapes.add_picture(
                    path, Inches(x), Inches(y), Inches(w), Inches(h)); return
            except Exception: pass
        s = self.slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = hex_to_rgb("EEF3F8")
        s.line.color.rgb = hex_to_rgb("AABBCC"); s.line.width = Pt(1.0)

    def _caption(self, text, x, y, w, h):
        self.add_textbox(x, y, w, h, text=f"▲  {text}",
            font_name=self.FONT_CN, font_size=11,
            color="555555", align=PP_ALIGN.CENTER)

    def _two_figs(self, points, figures, cx, top, cw, avail):
        """2张图：上方文字 + 下方左右并排两图"""
        CAP_H    = 0.20
        texts    = [str(p).strip() for p in points]
        txt_h    = avail * 0.35 if texts else 0
        img_zone_top = top + txt_h + (0.08 if texts else 0)
        img_zone_h   = self.FOOT_Y - img_zone_top - CAP_H

        if texts:
            fs, heights = _calc(texts, txt_h, cw)
            gap = max(0, txt_h - sum(heights)) / max(len(texts), 1)
            self._draw_texts(texts, cx, cw, top, fs, heights, gap)
            self.add_line(cx, top + txt_h + 0.02, cw, color="CCCCCC", width_pt=0.5)

        each_w = (cw - 0.12) / 2
        for i, fig in enumerate(figures[:2]):
            path = getattr(fig, "path", "") or ""
            cap  = getattr(fig, "caption", getattr(fig, "label", ""))
            asp  = _aspect(path) if (path and os.path.isfile(path)) else 1.5
            w, h = self._fit(each_w, img_zone_h, asp)
            x = cx + i * (each_w + 0.12) + (each_w - w) / 2
            y = img_zone_top + (img_zone_h - h) / 2
            self._draw_img(path, x, y, w, h)
            self._caption(cap, cx + i * (each_w + 0.12), y + h + 0.02, each_w, CAP_H)

    def _three_figs(self, points, figures, cx, top, cw, avail):
        """3张图：上方文字+1张大图，下方2张小图"""
        CAP_H  = 0.20
        texts  = [str(p).strip() for p in points]
        txt_h  = avail * 0.28 if texts else 0

        if texts:
            fs, heights = _calc(texts, txt_h, cw)
            gap = max(0, txt_h - sum(heights)) / max(len(texts), 1)
            self._draw_texts(texts, cx, cw, top, fs, heights, gap)
            self.add_line(cx, top + txt_h + 0.02, cw, color="CCCCCC", width_pt=0.5)

        zone_top = top + txt_h + (0.08 if texts else 0)
        zone_h   = self.FOOT_Y - zone_top - CAP_H
        row1_h   = zone_h * 0.55
        row2_h   = zone_h * 0.40

        # 第一行：1张大图居中
        fig0 = figures[0]
        path0 = getattr(fig0, "path", "") or ""
        cap0  = getattr(fig0, "caption", getattr(fig0, "label", ""))
        asp0  = _aspect(path0) if (path0 and os.path.isfile(path0)) else 1.5
        w0, h0 = self._fit(cw * 0.6, row1_h, asp0)
        x0 = cx + (cw - w0) / 2
        self._draw_img(path0, x0, zone_top, w0, h0)
        self._caption(cap0, cx, zone_top + h0 + 0.02, cw, CAP_H)

        # 第二行：2张小图
        row2_top = zone_top + row1_h + CAP_H + 0.05
        each_w   = (cw - 0.12) / 2
        for i, fig in enumerate(figures[1:3]):
            path = getattr(fig, "path", "") or ""
            cap  = getattr(fig, "caption", getattr(fig, "label", ""))
            asp  = _aspect(path) if (path and os.path.isfile(path)) else 1.5
            w, h = self._fit(each_w, row2_h, asp)
            x = cx + i * (each_w + 0.12) + (each_w - w) / 2
            self._draw_img(path, x, row2_top, w, h)
            self._caption(cap, cx + i * (each_w + 0.12), row2_top + h + 0.02, each_w, CAP_H)

    def _four_figs(self, points, figures, cx, top, cw, avail):
        """4张图：上方文字 + 下方2×2网格"""
        CAP_H  = 0.18
        texts  = [str(p).strip() for p in points]
        txt_h  = avail * 0.25 if texts else 0

        if texts:
            fs, heights = _calc(texts, txt_h, cw)
            gap = max(0, txt_h - sum(heights)) / max(len(texts), 1)
            self._draw_texts(texts, cx, cw, top, fs, heights, gap)
            self.add_line(cx, top + txt_h + 0.02, cw, color="CCCCCC", width_pt=0.5)

        zone_top = top + txt_h + (0.08 if texts else 0)
        zone_h   = self.FOOT_Y - zone_top - CAP_H * 2 - 0.05
        row_h    = (zone_h - 0.08) / 2
        each_w   = (cw - 0.12) / 2

        for idx, fig in enumerate(figures[:4]):
            row = idx // 2
            col = idx % 2
            path = getattr(fig, "path", "") or ""
            cap  = getattr(fig, "caption", getattr(fig, "label", ""))
            asp  = _aspect(path) if (path and os.path.isfile(path)) else 1.5
            w, h = self._fit(each_w, row_h, asp)
            x = cx + col * (each_w + 0.12) + (each_w - w) / 2
            y = zone_top + row * (row_h + CAP_H + 0.05) + (row_h - h) / 2
            self._draw_img(path, x, y, w, h)
            self._caption(cap, cx + col * (each_w + 0.12), y + h + 0.02, each_w, CAP_H)

    def _text_full(self, points, cx, cw, top, avail):
        if not points: return
        texts = [str(p).strip() for p in points]
        fs, heights = _calc(texts, avail, cw)
        gap = max(0, avail - sum(heights)) / max(len(texts), 1)
        self._draw_texts(texts, cx, cw, top, fs, heights, gap)

    def _rich(self, paragraph, text, size=14):
        for part in re.split(r"(\*\*.*?\*\*)", text):
            if part.startswith("**") and part.endswith("**"):
                r = paragraph.add_run(); r.text = part[2:-2]
                r.font.name = self.FONT_EN; r.font.size = Pt(size)
                r.font.bold = True; r.font.color.rgb = hex_to_rgb(self.COLOR_RED)
            elif part:
                r = paragraph.add_run(); r.text = part
                r.font.name = self.FONT_EN; r.font.size = Pt(size)
                r.font.color.rgb = hex_to_rgb(self.COLOR_BLACK)
