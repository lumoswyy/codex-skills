"""
PPT 生成器 - 禁用 AutoLayout 二次调整（会破坏精确布局）
"""

import os
from pptx import Presentation
from pptx.util import Inches
from typing import List
from ..models import (
    Page,
    PAGE_TYPE_COVER,
    PAGE_TYPE_TOC,
    PAGE_TYPE_SECTION,
    PAGE_TYPE_CONTENT_LIST,
    PAGE_TYPE_CONTENT_DETAIL,
    PAGE_TYPE_CONTENT_WITH_FIG,
    PAGE_TYPE_TIMELINE,
    PAGE_TYPE_ENDING,
)
from ..templates.cover import CoverTemplate
from ..templates.toc import TOCTemplate
from ..templates.section import SectionTemplate
from ..templates.content_list import ContentListTemplate
from ..templates.content_detail import ContentDetailTemplate
from ..templates.content_figure import ContentFigureTemplate
from ..templates.timeline import TimelineTemplate
from ..templates.ending import EndingTemplate
from ..config import default_config

NO_PAGE_NUM = {PAGE_TYPE_COVER, PAGE_TYPE_ENDING}


class PPTXGenerator:
    def __init__(self, config=None):
        self.config = config or default_config

    def generate(self, pages: List[Page], output_path: str = "output.pptx") -> str:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        sections = self._extract_sections(pages)
        total_pages = sum(1 for p in pages if p.page_type not in NO_PAGE_NUM)
        cur_page = 0

        for page in pages:
            slide = prs.slides.add_slide(blank)
            if page.page_type not in NO_PAGE_NUM:
                cur_page += 1
                pn = cur_page
            else:
                pn = None
            self._render_page(slide, page, sections, pn, total_pages)
            # ✅ 不再调用 AutoLayout，避免破坏精确布局

        # 生成后质量检查
        warnings = self._quality_check(pages)
        for w in warnings:
            self.config.logger.warning(w)

        prs.save(output_path)
        return output_path

    def _quality_check(self, pages: List[Page]) -> List[str]:
        """生成后质量自动检查，返回警告列表"""
        warnings = []
        if len(pages) <= 2:
            warnings.append(f"PPT 只有 {len(pages)} 页，可能解析失败")

        for i, page in enumerate(pages):
            # 检查内容页是否有内容
            if page.page_type == PAGE_TYPE_CONTENT_LIST:
                items = getattr(page.data, "items", []) or []
                title = getattr(page.data, "title", "未知")
                if len(items) == 0:
                    warnings.append(f"第{i + 1}页 ({title}) 没有内容项，可能为空白页")

            # 检查图文页的图片路径是否存在
            if page.page_type == PAGE_TYPE_CONTENT_WITH_FIG:
                figs = getattr(page.data, "figures", []) or []
                title = getattr(page.data, "title", "未知")
                if not figs:
                    fig = getattr(page.data, "figure", None)
                    if fig:
                        figs = [fig]
                for fig in figs:
                    fig_path = getattr(fig, "path", "") or ""
                    if fig_path and not os.path.isfile(fig_path):
                        warnings.append(f"第{i + 1}页 ({title}) 图片不存在: {fig_path}")

        return warnings

    def _extract_sections(self, pages):
        seen, result = set(), []
        for p in pages:
            if p.page_type == PAGE_TYPE_SECTION:
                t = getattr(p.data, "part_title", "") or ""
                if t and t not in seen:
                    result.append(t)
                    seen.add(t)
        return result

    def _get_current_section(self, page, sections):
        if page.page_type == PAGE_TYPE_SECTION:
            return getattr(page.data, "part_title", "") or ""
        # 优先用数据里显式绑定的章节
        explicit = getattr(page.data, "current_section", "") or ""
        if explicit and explicit in sections:
            return explicit
        # fallback：字符串匹配
        title = getattr(page.data, "title", "") or ""
        for sec in sections:
            if sec in title or title in sec:
                return sec
        return sections[0] if sections else ""

    def _render_page(self, slide, page, sections, page_num, total_pages):
        pt = page.page_type
        data = page.data
        cur = self._get_current_section(page, sections)
        kw = dict(
            sections=sections,
            current_section=cur,
            page_num=page_num,
            total_pages=total_pages,
        )

        if pt == PAGE_TYPE_COVER:
            CoverTemplate(slide, self.config).render(data)
        elif pt == PAGE_TYPE_TOC:
            t = TOCTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()
        elif pt == PAGE_TYPE_SECTION:
            t = SectionTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()
        elif pt == PAGE_TYPE_CONTENT_LIST:
            t = ContentListTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()
        elif pt == PAGE_TYPE_CONTENT_DETAIL:
            t = ContentDetailTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()
        elif pt == PAGE_TYPE_CONTENT_WITH_FIG:
            t = ContentFigureTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()
        elif pt == PAGE_TYPE_TIMELINE:
            t = TimelineTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()
        elif pt == PAGE_TYPE_ENDING:
            EndingTemplate(slide, self.config).render(data)
        else:
            t = ContentListTemplate(slide, self.config, **kw)
            t.render(data)
            t.draw_page_number()


def generate_ppt(pages, output_path="output.pptx"):
    return PPTXGenerator().generate(pages, output_path)
