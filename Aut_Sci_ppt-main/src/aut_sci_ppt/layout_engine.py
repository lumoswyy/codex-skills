"""
自适应布局引擎
- 根据内容量自动调整字号、行间距、内容区高度
- 确保每页不溢出、重点突出、排版专业
"""
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List, Tuple


# ── 字号自适应规则 ────────────────────────────────────────────

def calc_font_size(text: str, box_width_inch: float, box_height_inch: float,
                   base_size: int = 16, min_size: int = 10, max_size: int = 28) -> int:
    """
    根据文字量和文本框尺寸自动计算字号
    估算每行字符数、行数，确保内容不溢出
    """
    if not text:
        return base_size

    # 每个字符平均宽度（英寸）：中文约 0.18，英文约 0.10（以 16pt 为基准）
    char_w_cn = 0.155   # 16pt 中文字符宽
    chars_per_line = int(box_width_inch / char_w_cn)
    if chars_per_line < 1:
        chars_per_line = 1

    line_h_inch = 0.28  # 16pt 行高（含行间距）
    max_lines = int(box_height_inch / line_h_inch)
    if max_lines < 1:
        max_lines = 1

    text_len = len(text)
    needed_lines = max(1, -(-text_len // chars_per_line))  # ceil division

    if needed_lines <= max_lines:
        # 内容放得下，尝试适当放大
        scale = min(max_lines / max(needed_lines, 1), 1.3)
        size = min(int(base_size * scale), max_size)
    else:
        # 内容放不下，缩小字号
        scale = max_lines / needed_lines
        size = max(int(base_size * scale), min_size)

    return size


def calc_list_layout(items: list, avail_height: float,
                     base_font: int = 15, min_font: int = 10) -> Tuple[int, float]:
    """
    根据列表项数量，计算最优字号和每项高度
    返回 (font_size, item_height_inch)
    """
    n = max(len(items), 1)

    # 理想每项高度（含行间距）
    ideal_h = avail_height / n

    # 字号与行高的映射关系
    SIZE_TO_H = {
        15: 0.62, 14: 0.57, 13: 0.52, 12: 0.47,
        11: 0.43, 10: 0.39
    }

    font_size = base_font
    for size in range(base_font, min_font - 1, -1):
        h = SIZE_TO_H.get(size, size * 0.038)
        if h <= ideal_h:
            font_size = size
            break

    item_h = SIZE_TO_H.get(font_size, font_size * 0.038)
    # 如果内容少，适当增大行高（但不超过 0.85）
    if n <= 4:
        item_h = min(ideal_h * 0.85, 0.85)

    return font_size, item_h


def calc_title_font(title: str, box_width: float,
                    base: int = 22, min_s: int = 14, max_s: int = 28) -> int:
    """页眉标题字号自适应"""
    if not title:
        return base
    # 估算：宽屏下每英寸约放 3 个中文字符（22pt）
    capacity = int(box_width * 2.8)
    if len(title) <= capacity:
        return min(base + 2, max_s)
    elif len(title) <= capacity * 1.5:
        return base
    else:
        return max(min_s, base - 4)


# ── 文本框自适应调整器 ────────────────────────────────────────

class AutoLayout:
    """
    在页面渲染完成后，对文本框进行二次自适应调整
    主要处理：字号、行间距、对齐方式
    """

    @staticmethod
    def adjust_textbox(txBox, max_width_inch: float = None,
                       max_height_inch: float = None,
                       min_font: int = 10):
        """
        自动调整文本框内所有 Run 的字号，确保不溢出
        """
        tf = txBox.text_frame
        # 获取文本框实际尺寸
        w = txBox.width / 914400   # EMU -> inch
        h = txBox.height / 914400

        if max_width_inch:
            w = min(w, max_width_inch)
        if max_height_inch:
            h = min(h, max_height_inch)

        full_text = tf.text
        if not full_text.strip():
            return

        # 计算当前最大字号
        cur_max = 0
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    cur_max = max(cur_max, run.font.size / 12700)

        if cur_max == 0:
            return

        optimal = calc_font_size(full_text, w, h,
                                 base_size=int(cur_max),
                                 min_size=min_font,
                                 max_size=int(cur_max) + 4)

        if optimal == cur_max:
            return

        ratio = optimal / cur_max
        for para in tf.paragraphs:
            # 自动调整行间距
            from pptx.util import Pt
            para.line_spacing = Pt(optimal * 1.35)
            for run in para.runs:
                if run.font.size:
                    new_size = max(int(run.font.size / 12700 * ratio), min_font)
                    run.font.size = Pt(new_size)

    @staticmethod
    def adjust_slide(slide, nav_w: float = 1.75, min_font: int = 10):
        """
        对整张幻灯片的所有文本框进行自适应调整
        跳过导航栏区域（x < nav_w）
        """
        from pptx.shapes.autoshape import Shape
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # 跳过导航栏内的形状
            left_inch = shape.left / 914400
            if left_inch < nav_w - 0.1:
                continue
            w = shape.width / 914400
            h = shape.height / 914400
            # 跳过太小的文本框（标题栏等固定尺寸）
            if h < 0.3 or w < 0.5:
                continue
            AutoLayout.adjust_textbox(shape, max_width_inch=w, max_height_inch=h,
                                      min_font=min_font)


# ── 内容精简规则 ──────────────────────────────────────────────

def truncate_text(text: str, max_chars: int = 80) -> str:
    """超长文本自动截断并加省略号"""
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 1] + "…"


def split_long_items(items: list, max_chars: int = 60) -> list:
    """
    超长列表项自动折行处理
    若单项超出 max_chars，拆分为主句 + 补充说明
    """
    result = []
    for item in items:
        text = item.text if hasattr(item, "text") else str(item)
        if len(text) > max_chars:
            # 尝试在标点处断开
            import re
            parts = re.split(r"[，,；;]", text, maxsplit=1)
            if len(parts) == 2 and len(parts[0]) < max_chars:
                # 保留主句，副句作为子项
                if hasattr(item, "text"):
                    item.text = parts[0]
                result.append(item)
                # 补充说明作为缩进子项（加前缀空格）
                from aut_sci_ppt.models import ListItem
                result.append(ListItem(icon=" ", text=f"  {parts[1]}"))
                continue
        result.append(item)
    return result
