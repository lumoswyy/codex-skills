"""封面页 - 信息区均匀分布版"""
from .base import BaseTemplate, hex_to_rgb
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class CoverTemplate(BaseTemplate):
    def render(self, data):
        self.add_rect(0, 0, 0.5, self.H, fill_color=self.COLOR_PRIMARY)
        self.add_rect(0.5, 0, self.W - 0.5, 0.07, fill_color=self.COLOR_GOLD)
        self.add_rect(0, 6.72, self.W, 0.78, fill_color=self.COLOR_PRIMARY)
        self.add_rect(0, 6.72, self.W, 0.07, fill_color=self.COLOR_GOLD)

        title = getattr(data, "title", "") or "科研汇报"
        fs = 44 if len(title) <= 16 else (36 if len(title) <= 26 else 28)
        tb = self.slide.shapes.add_textbox(
            Inches(1.2), Inches(1.7), Inches(11.0), Inches(1.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = title
        run.font.name = self.FONT_CN; run.font.size = Pt(fs)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(self.COLOR_PRIMARY)

        self.add_line(2.2, 3.65, 9.0, color=self.COLOR_GOLD, width_pt=2.0)

        LABEL_DEFS = [
            ("汇  报  人", "author"),
            ("指导教师",   "advisor"),
            ("研究方向",   "direction"),
            ("日      期", "date"),
        ]
        fields = [(lbl, getattr(data, attr, ""))
                  for lbl, attr in LABEL_DEFS if getattr(data, attr, "")]
        if not fields: return

        FONT_SIZE  = 20
        SEP_W      = 0.28
        LABEL_W    = 1.65
        VALUE_W    = 5.5
        TOTAL_W    = LABEL_W + SEP_W + VALUE_W
        block_left = (self.W - TOTAL_W) / 2

        # 信息区可用范围：分割线下 3.75 到底部色条上 6.72
        INFO_TOP  = 3.78
        INFO_BOT  = 6.55
        INFO_H    = INFO_BOT - INFO_TOP
        n         = len(fields)
        # 均匀分布：每项高度 = 总高 / 项数
        LINE_H    = INFO_H / n

        for i, (label, value) in enumerate(fields):
            row_top = INFO_TOP + i * LINE_H
            # 文字垂直居中在 LINE_H 内：顶部偏移 = (LINE_H - 单行高) / 2
            text_offset = max(0, (LINE_H - 0.36) / 2)

            # 标签
            tb_l = self.slide.shapes.add_textbox(
                Inches(block_left), Inches(row_top + text_offset),
                Inches(LABEL_W), Inches(0.44))
            tf_l = tb_l.text_frame; tf_l.word_wrap = False
            p_l = tf_l.paragraphs[0]; p_l.alignment = PP_ALIGN.DISTRIBUTE
            r = p_l.add_run(); r.text = label
            r.font.name = self.FONT_KAI; r.font.size = Pt(FONT_SIZE)
            r.font.bold = True
            r.font.color.rgb = hex_to_rgb(self.COLOR_PRIMARY)

            # 冒号
            tb_s = self.slide.shapes.add_textbox(
                Inches(block_left + LABEL_W), Inches(row_top + text_offset),
                Inches(SEP_W), Inches(0.44))
            p_s = tb_s.text_frame.paragraphs[0]; p_s.alignment = PP_ALIGN.CENTER
            r_s = p_s.add_run(); r_s.text = "："
            r_s.font.name = self.FONT_KAI; r_s.font.size = Pt(FONT_SIZE)
            r_s.font.bold = True
            r_s.font.color.rgb = hex_to_rgb(self.COLOR_PRIMARY)

            # 值（允许换行，给足高度）
            tb_v = self.slide.shapes.add_textbox(
                Inches(block_left + LABEL_W + SEP_W), Inches(row_top + text_offset),
                Inches(VALUE_W), Inches(LINE_H - text_offset))
            tf_v = tb_v.text_frame; tf_v.word_wrap = True
            p_v = tf_v.paragraphs[0]; p_v.alignment = PP_ALIGN.CENTER
            r_v = p_v.add_run(); r_v.text = value
            r_v.font.name = self.FONT_KAI; r_v.font.size = Pt(FONT_SIZE)
            r_v.font.bold = True
            r_v.font.color.rgb = hex_to_rgb(self.COLOR_BLACK)
