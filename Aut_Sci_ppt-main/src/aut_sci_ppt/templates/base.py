"""
基础模板 - 导航栏自适应字号 彻底修复版
颜色/字体优先从 config 读取，类常量作为默认兜底
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import List


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class BaseTemplate:
    # 类常量：作为 config 未提供时的兜底默认值
    COLOR_PRIMARY = "1E3A5F"
    COLOR_RED = "BF0000"
    COLOR_BLACK = "000000"
    COLOR_WHITE = "FFFFFF"
    COLOR_GOLD = "E8B339"
    COLOR_NAV_INACTIVE = "9BBDD8"
    COLOR_NAV_BG = "122540"

    FONT_CN = "微软雅黑"
    FONT_KAI = "楷体"
    FONT_EN = "Times New Roman"

    NAV_W = 2.0
    CONTENT_X = 2.15

    def __init__(
        self,
        slide,
        config,
        sections=None,
        current_section="",
        page_num=None,
        total_pages=None,
    ):
        self.slide = slide
        self.config = config
        self.W = 13.33
        self.H = 7.5
        self.sections = sections or []
        self.current_section = current_section
        self.page_num = page_num
        self.total_pages = total_pages

        # 从 config 同步颜色（让 config.py 的配置真正生效）
        # config.DEFAULT_COLORS 的 key → 对应的实例属性
        if config is not None:
            colors = getattr(config, "colors", None) or getattr(
                config, "DEFAULT_COLORS", {}
            )
            if colors:
                self.COLOR_PRIMARY = colors.get("primary", self.COLOR_PRIMARY)
                self.COLOR_RED = colors.get("accent_red", self.COLOR_RED)
                self.COLOR_GOLD = colors.get("accent_gold", self.COLOR_GOLD)
                self.COLOR_BLACK = colors.get("text_black", self.COLOR_BLACK)
                self.COLOR_WHITE = colors.get("title_bar_text", self.COLOR_WHITE)

            fonts = getattr(config, "fonts", None) or getattr(
                config, "DEFAULT_FONTS", {}
            )
            if fonts:
                cover_info = fonts.get("cover_info", {})
                if cover_info.get("name"):
                    self.FONT_KAI = cover_info["name"]  # 楷体（封面信息区用）
                body = fonts.get("body", {})
                if body.get("name"):
                    self.FONT_EN = body["name"]  # Times New Roman（正文用）

    def add_rect(self, left, top, width, height, fill_color=None, line_color=None):
        shape = self.slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        fill = shape.fill
        if fill_color:
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(fill_color)
        else:
            fill.background()
        line = shape.line
        if line_color:
            line.color.rgb = hex_to_rgb(line_color)
            line.width = Pt(0.75)
        else:
            line.fill.background()
        return shape

    def add_textbox(
        self,
        left,
        top,
        width,
        height,
        text="",
        font_name=None,
        font_size=18,
        bold=False,
        color="000000",
        align=PP_ALIGN.LEFT,
        wrap=True,
    ):
        tb = self.slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name or self.FONT_CN
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)
        return tb

    def add_line(self, left, top, width, color="1E3A5F", width_pt=1.5):
        c = self.slide.shapes.add_connector(
            1, Inches(left), Inches(top), Inches(left + width), Inches(top)
        )
        c.line.color.rgb = hex_to_rgb(color)
        c.line.width = Pt(width_pt)
        return c

    def draw_page_number(self):
        if self.page_num is None:
            return
        txt = (
            f"{self.page_num} / {self.total_pages}"
            if self.total_pages
            else str(self.page_num)
        )
        self.add_textbox(
            self.W - 1.1,
            self.H - 0.36,
            0.95,
            0.3,
            text=txt,
            font_name=self.FONT_EN,
            font_size=11,
            bold=False,
            color="888888",
            align=PP_ALIGN.RIGHT,
        )

    def draw_nav_sidebar(self):
        self.add_rect(0, 0, self.NAV_W, self.H, fill_color=self.COLOR_NAV_BG)
        if not self.sections:
            return

        n = len(self.sections)

        # ── 滑动窗口：最多显示5个章节，以当前章节居中 ──────
        WINDOW = 5
        if n <= WINDOW:
            visible_sections = self.sections
            visible_offset = 0
        else:
            try:
                cur_idx = self.sections.index(self.current_section)
            except ValueError:
                cur_idx = 0
            half = WINDOW // 2
            start = cur_idx - half
            start = max(0, min(start, n - WINDOW))
            visible_sections = self.sections[start : start + WINDOW]
            visible_offset = start

        # Logo 区
        LOGO_H = 0.72
        self.add_rect(0, 0, self.NAV_W, LOGO_H, fill_color="0D1E35")
        self.add_textbox(
            0.05,
            0.07,
            self.NAV_W - 0.1,
            0.35,
            text="RESEARCH",
            font_name=self.FONT_EN,
            font_size=13,
            bold=True,
            color=self.COLOR_GOLD,
            align=PP_ALIGN.CENTER,
        )
        self.add_textbox(
            0.05,
            0.40,
            self.NAV_W - 0.1,
            0.24,
            text="NAVIGATOR",
            font_name=self.FONT_EN,
            font_size=8,
            bold=False,
            color="5A8AB0",
            align=PP_ALIGN.CENTER,
        )
        self.add_line(
            0.1, LOGO_H, self.NAV_W - 0.2, color=self.COLOR_GOLD, width_pt=0.8
        )

        BOTTOM_PAD = 0.28
        n_visible = len(visible_sections)
        avail_h = self.H - LOGO_H - BOTTOM_PAD
        item_h = avail_h / n_visible  # 按实际显示数等高

        TITLE_H = item_h * 0.58
        max_len = max(len(s) for s in visible_sections)
        chars_per_line = max(1, int(1.82 * 72 / 14 / 1.1))

        if max_len <= chars_per_line:
            title_pt = int(TITLE_H * 72 / 1.45)
        else:
            title_pt = int(TITLE_H * 72 / 2.9)

        title_pt = max(9, min(title_pt, 20))
        num_pt = max(8, min(int(title_pt * 0.7), 13))

        PAD_X = 0.12

        for i, sec_title in enumerate(visible_sections):
            top = LOGO_H + i * item_h
            is_cur = sec_title == self.current_section

            if is_cur:
                self.add_rect(0, top, self.NAV_W, item_h, fill_color="1C3658")
                self.add_rect(0, top, 0.07, item_h, fill_color=self.COLOR_GOLD)
                n_color, t_color, t_bold = self.COLOR_GOLD, self.COLOR_WHITE, True
            else:
                n_color, t_color, t_bold = "3D6A99", self.COLOR_NAV_INACTIVE, False

            # 编号（item_h 上方 30%）
            num_h = item_h * 0.30
            self.add_textbox(
                PAD_X,
                top + item_h * 0.03,
                self.NAV_W - PAD_X - 0.04,
                num_h,
                text=f"{visible_offset + i + 1:02d}",
                font_name=self.FONT_EN,
                font_size=num_pt,
                bold=is_cur,
                color=n_color,
            )

            # 标题（剩余 60%，开启换行）
            title_top = top + num_h + item_h * 0.04
            title_box_h = item_h * 0.62
            txBox = self.slide.shapes.add_textbox(
                Inches(PAD_X),
                Inches(title_top),
                Inches(self.NAV_W - PAD_X - 0.04),
                Inches(title_box_h),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = sec_title
            run.font.name = self.FONT_CN
            run.font.size = Pt(title_pt)
            run.font.bold = t_bold
            run.font.color.rgb = hex_to_rgb(t_color)

            # 分隔线
            if i < n_visible - 1 and not is_cur:
                self.add_line(
                    PAD_X,
                    top + item_h - 0.008,
                    self.NAV_W - PAD_X * 2,
                    color="1A3358",
                    width_pt=0.3,
                )

        self.add_line(
            0.1,
            self.H - BOTTOM_PAD + 0.03,
            self.NAV_W - 0.2,
            color=self.COLOR_GOLD,
            width_pt=0.5,
        )

    def draw_header_bar(self, part_num="", title=""):
        cw = self.W - self.CONTENT_X
        self.add_rect(self.CONTENT_X, 0, cw, 0.62, fill_color=self.COLOR_PRIMARY)
        self.add_rect(self.CONTENT_X, 0, 0.07, 0.62, fill_color=self.COLOR_GOLD)
        x = self.CONTENT_X + 0.18
        if part_num:
            self.add_textbox(
                x,
                0.06,
                0.65,
                0.5,
                text=str(part_num),
                font_name=self.FONT_EN,
                font_size=20,
                bold=True,
                color=self.COLOR_GOLD,
            )
            x += 0.6
        if title:
            from ..layout_engine import calc_title_font

            ts = calc_title_font(title, self.W - x - 0.2, base=20, min_s=13, max_s=22)
            self.add_textbox(
                x,
                0.06,
                self.W - x - 0.2,
                0.5,
                text=title,
                font_name=self.FONT_CN,
                font_size=ts,
                bold=True,
                color=self.COLOR_WHITE,
            )

    def render(self, data):
        raise NotImplementedError
