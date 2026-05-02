"""
PPT Agent 模板系统 - 带图片的详情内容页模板
"""
import os
from .base import BaseTemplate


class ContentDetailImageTemplate(BaseTemplate):
    """带图片的详情内容页模板"""
    
    def render(self, slide, data):
        """渲染带图片的详情内容页"""
        # 获取布局配置
        layout = getattr(data, 'layout', 'right')  # 默认右图左文
        
        if layout == 'left':
            self._render_left_image(slide, data)
        elif layout == 'right':
            self._render_right_image(slide, data)
        elif layout == 'top':
            self._render_top_image(slide, data)
        elif layout == 'bottom':
            self._render_bottom_image(slide, data)
        else:
            # 默认右图左文
            self._render_right_image(slide, data)
    
    def _render_right_image(self, slide, data):
        """右图左文布局"""
        # 页面标题
        if hasattr(data, 'title') and data.title:
            self._add_title(
                slide,
                data.title,
                font_type='title',
                top=0.3
            )
        
        # 副标题
        if hasattr(data, 'subtitle') and data.subtitle:
            self._add_text_box(
                slide,
                data.subtitle,
                font_type='subtitle',
                top=1.1,
                alignment='center'
            )
        
        # 左侧文字内容
        content_top = 1.8
        
        # 研究背景
        if hasattr(data, 'background') and data.background:
            self._add_text_box(
                slide,
                f"研究背景：{data.background}",
                font_type='body',
                top=content_top,
                left=0.5,
                width=4.5,
                height=1.0
            )
            content_top += 1.2
        
        # 研究要点
        if hasattr(data, 'points') and data.points:
            points_text = "研究要点：\n" + "\n".join([f"{i+1}. {p}" for i, p in enumerate(data.points)])
            self._add_text_box(
                slide,
                points_text,
                font_type='body',
                top=content_top,
                left=0.5,
                width=4.5,
                height=2.5
            )
            content_top += 2.5
        
        # 研究成果
        if hasattr(data, 'results') and data.results:
            results_text = "研究成果：\n" + "\n".join([f"• {r}" for r in data.results])
            self._add_text_box(
                slide,
                results_text,
                font_type='body',
                top=content_top,
                left=0.5,
                width=4.5,
                height=1.5
            )
        
        # 右侧图片
        if hasattr(data, 'images') and data.images:
            for img_item in data.images:
                if os.path.exists(img_item.path):
                    self._add_image(
                        slide,
                        img_item.path,
                        position='right',
                        size=img_item.size,
                        caption=img_item.caption
                    )
    
    def _render_left_image(self, slide, data):
        """左图右文布局"""
        # 页面标题
        if hasattr(data, 'title') and data.title:
            self._add_title(
                slide,
                data.title,
                font_type='title',
                top=0.3
            )
        
        # 副标题
        if hasattr(data, 'subtitle') and data.subtitle:
            self._add_text_box(
                slide,
                data.subtitle,
                font_type='subtitle',
                top=1.1,
                alignment='center'
            )
        
        # 左侧图片
        if hasattr(data, 'images') and data.images:
            for img_item in data.images:
                if os.path.exists(img_item.path):
                    self._add_image(
                        slide,
                        img_item.path,
                        position='left',
                        size=img_item.size,
                        caption=img_item.caption
                    )
        
        # 右侧文字内容
        content_top = 1.8
        
        # 研究背景
        if hasattr(data, 'background') and data.background:
            self._add_text_box(
                slide,
                f"研究背景：{data.background}",
                font_type='body',
                top=content_top,
                left=5.2,
                width=4.5,
                height=1.0
            )
            content_top += 1.2
        
        # 研究要点
        if hasattr(data, 'points') and data.points:
            points_text = "研究要点：\n" + "\n".join([f"{i+1}. {p}" for i, p in enumerate(data.points)])
            self._add_text_box(
                slide,
                points_text,
                font_type='body',
                top=content_top,
                left=5.2,
                width=4.5,
                height=2.5
            )
            content_top += 2.5
        
        # 研究成果
        if hasattr(data, 'results') and data.results:
            results_text = "研究成果：\n" + "\n".join([f"• {r}" for r in data.results])
            self._add_text_box(
                slide,
                results_text,
                font_type='body',
                top=content_top,
                left=5.2,
                width=4.5,
                height=1.5
            )
    
    def _render_top_image(self, slide, data):
        """上图下文布局"""
        # 顶部图片
        if hasattr(data, 'images') and data.images:
            for img_item in data.images:
                if os.path.exists(img_item.path):
                    self._add_image(
                        slide,
                        img_item.path,
                        position='top',
                        size=img_item.size,
                        caption=img_item.caption
                    )
        
        # 页面标题
        if hasattr(data, 'title') and data.title:
            self._add_title(
                slide,
                data.title,
                font_type='title',
                top=3.5
            )
        
        # 文字内容
        content_top = 4.2
        
        # 研究要点
        if hasattr(data, 'points') and data.points:
            points_text = "研究要点：\n" + "\n".join([f"{i+1}. {p}" for i, p in enumerate(data.points)])
            self._add_text_box(
                slide,
                points_text,
                font_type='body',
                top=content_top,
                height=2.8
            )
    
    def _render_bottom_image(self, slide, data):
        """下图上文布局"""
        # 页面标题
        if hasattr(data, 'title') and data.title:
            self._add_title(
                slide,
                data.title,
                font_type='title',
                top=0.3
            )
        
        # 文字内容
        content_top = 1.0
        
        # 研究要点
        if hasattr(data, 'points') and data.points:
            points_text = "研究要点：\n" + "\n".join([f"{i+1}. {p}" for i, p in enumerate(data.points)])
            self._add_text_box(
                slide,
                points_text,
                font_type='body',
                top=content_top,
                height=2.8
            )
        
        # 底部图片
        if hasattr(data, 'images') and data.images:
            for img_item in data.images:
                if os.path.exists(img_item.path):
                    self._add_image(
                        slide,
                        img_item.path,
                        position='bottom',
                        size=img_item.size,
                        caption=img_item.caption
                    )
    
    def _add_image(self, slide, image_path, position='right', size=None, caption=''):
        """添加图片"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # 默认尺寸
        if size is None:
            size = {'width': 4, 'height': 3}
        
        img_width = Inches(size.get('width', 4))
        img_height = Inches(size.get('height', 3))
        
        # 根据位置计算坐标
        if position == 'right':
            left = Inches(5.5)
            top = Inches(2.0)
        elif position == 'left':
            left = Inches(0.5)
            top = Inches(2.0)
        elif position == 'top':
            left = Inches(1.5)
            top = Inches(0.5)
        elif position == 'bottom':
            left = Inches(1.5)
            top = Inches(4.0)
        else:
            left = Inches(5.5)
            top = Inches(2.0)
        
        try:
            # 添加图片
            slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)
            
            # 添加图注
            if caption:
                caption_top = top + img_height + Inches(0.1)
                txBox = slide.shapes.add_textbox(
                    left, caption_top, img_width, Inches(0.3)
                )
                text_frame = txBox.text_frame
                text_frame.word_wrap = True
                
                p = text_frame.paragraphs[0]
                p.text = caption
                p.font.name = self._get_font_config('footer')['name']
                p.font.size = Pt(12)
                p.font.color.rgb = self._get_color('text_light')
                p.alignment = PP_ALIGN.CENTER
        except Exception as e:
            # 图片加载失败，静默跳过
            pass
    
    def _add_title(self, slide, text, font_type='title', top=0.3, font_config=None):
        """添加主标题"""
        if font_config is None:
            font_config = self._get_font_config(font_type)
        
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        width = self.page_width - 1
        
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(top), Inches(width), Inches(1)
        )
        text_frame = txBox.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.name = font_config['name']
        p.font.size = Pt(font_config['size'])
        p.font.bold = font_config.get('bold', False)
        p.font.color.rgb = self._get_color('primary')
        p.alignment = PP_ALIGN.LEFT
        
        return txBox
